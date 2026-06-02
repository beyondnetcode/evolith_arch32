import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side

out_dir = "/Users/beyondnet/Source/evolith/reference/governance/sdlc/assets"
os.makedirs(out_dir, exist_ok=True)

# Colors
DARK_BG = RGBColor(18, 18, 18)        # #121212 Onyx/Black
GOLD_ACCENT = RGBColor(202, 197, 178) # #CAC5B2 Golden/Sand
WHITE_TEXT = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(160, 160, 160)

# --- PPTX HELPERS ---
def apply_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title(slide, text):
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = text
        for p in title_shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = GOLD_ACCENT
                r.font.size = Pt(36)
                r.font.bold = True
            p.alignment = PP_ALIGN.LEFT
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.5)
        title_shape.width = Inches(9.0)
        title_shape.height = Inches(1.0)

def add_conceptual_box(slide, x, y, w, h, text, color=GOLD_ACCENT, bg_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill = shape.fill
    if bg_color:
        fill.solid()
        fill.fore_color.rgb = bg_color
    else:
        fill.background()
    line = shape.line
    line.color.rgb = color
    line.width = Pt(2)
    
    tf = shape.text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.color.rgb = color if not bg_color else DARK_BG
            r.font.size = Pt(14)
            r.font.bold = True

def add_arrow(slide, x, y, w, h, color=GOLD_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    line = shape.line
    line.color.rgb = color

def create_deck(filename, title, subtitle):
    prs = Presentation()
    
    # Portada
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_theme(slide)
    
    add_title(slide, title)
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.color.rgb = WHITE_TEXT
    p.font.size = Pt(24)
    
    # Logo / Icon conceptual
    add_conceptual_box(slide, 4.0, 4.0, 2.0, 2.0, "EVOLITH\nCORE", GOLD_ACCENT, DARK_BG)
    
    return prs

def save_deck(prs, filename):
    prs.save(os.path.join(out_dir, filename))

# --- DECK 1: PROPUESTA ---
def build_deck_1():
    prs = create_deck("evolith-sdlc-value-proposition.pptx", "Evolith: Gobernanza Inteligente del SDLC", "Transforme la entrega de tecnología en valor de negocio.")
    
    # Slide: Flujo SDLC (Diagrama de Cajas)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank
    apply_dark_theme(slide)
    add_title(slide, "Comparativa SDLC: Tradicional vs Evolith")
    
    # Diagrama Tradicional (Caos)
    add_conceptual_box(slide, 0.5, 2.0, 2.0, 1.0, "Silos de\nNegocio", RGBColor(150,0,0))
    add_conceptual_box(slide, 4.0, 2.0, 2.0, 1.0, "Desarrollo\nAislado", RGBColor(150,0,0))
    add_conceptual_box(slide, 7.5, 2.0, 2.0, 1.0, "QA\nTardío", RGBColor(150,0,0))
    
    # Diagrama Evolith (Ordenado)
    add_conceptual_box(slide, 0.5, 4.5, 9.0, 2.0, "Gobernanza SDLC Evolith (End-to-End Trazabilidad)", GOLD_ACCENT)
    add_conceptual_box(slide, 1.0, 5.0, 1.5, 1.0, "Ideación", DARK_BG, GOLD_ACCENT)
    add_arrow(slide, 2.6, 5.3, 0.4, 0.4)
    add_conceptual_box(slide, 3.1, 5.0, 1.5, 1.0, "Construcción", DARK_BG, GOLD_ACCENT)
    add_arrow(slide, 4.7, 5.3, 0.4, 0.4)
    add_conceptual_box(slide, 5.2, 5.0, 1.5, 1.0, "Quality Gates", DARK_BG, GOLD_ACCENT)
    add_arrow(slide, 6.8, 5.3, 0.4, 0.4)
    add_conceptual_box(slide, 7.3, 5.0, 1.5, 1.0, "Producción", DARK_BG, GOLD_ACCENT)
    
    # Slide: Quality Gates (Swimlanes)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_theme(slide)
    add_title(slide, "Carriles de Gobernanza y Puertas de Decisión")
    
    # Carriles
    add_conceptual_box(slide, 0.5, 2.0, 9.0, 1.2, "Negocio (Sponsor, PO)", WHITE_TEXT)
    add_conceptual_box(slide, 0.5, 3.5, 9.0, 1.2, "Ingeniería (Arquitectura, Devs)", WHITE_TEXT)
    add_conceptual_box(slide, 0.5, 5.0, 9.0, 1.2, "Operaciones (QA, Release)", WHITE_TEXT)
    
    # Nodos
    add_conceptual_box(slide, 2.0, 2.2, 1.0, 0.8, "Gate 1", DARK_BG, GOLD_ACCENT)
    add_conceptual_box(slide, 5.0, 3.7, 1.0, 0.8, "Gate 2", DARK_BG, GOLD_ACCENT)
    add_conceptual_box(slide, 8.0, 5.2, 1.0, 0.8, "Gate 3", DARK_BG, GOLD_ACCENT)
    
    save_deck(prs, "evolith-sdlc-value-proposition.pptx")

# --- DECK 2: CASO UMS ---
def build_deck_2():
    prs = create_deck("evolith-ums-case-study.pptx", "UMS: De la Incertidumbre a la Entrega Continua", "Caso Práctico: Portal de Autoservicio Industrial")
    
    # Slide: Causa-Efecto (Ishikawa conceptual)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_theme(slide)
    add_title(slide, "La Raíz de los Retrasos (Antes de Evolith)")
    
    add_arrow(slide, 1.0, 4.0, 7.0, 0.2)
    add_conceptual_box(slide, 8.0, 3.5, 1.5, 1.0, "Sobrecostes\nRetrasos", RGBColor(150,0,0))
    add_conceptual_box(slide, 2.0, 2.0, 2.0, 1.0, "Silos Org.", GRAY_TEXT)
    add_conceptual_box(slide, 4.0, 5.0, 2.0, 1.0, "Deuda Técnica", GRAY_TEXT)
    add_conceptual_box(slide, 6.0, 2.0, 2.0, 1.0, "Pruebas Tardías", GRAY_TEXT)
    
    # Slide: Métricas Antes vs Después
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_theme(slide)
    add_title(slide, "El Retorno de Inversión (Resultados)")
    
    # Gráficos de barra conceptuales
    add_conceptual_box(slide, 1.0, 3.0, 2.0, 4.0, "TTM Anterior\n(18 meses)", GRAY_TEXT)
    add_conceptual_box(slide, 3.5, 5.0, 2.0, 2.0, "TTM Evolith\n(8 meses)", DARK_BG, GOLD_ACCENT)
    add_conceptual_box(slide, 6.0, 2.0, 3.0, 1.0, "100% Trazabilidad Lograda", DARK_BG, GOLD_ACCENT)
    
    save_deck(prs, "evolith-ums-case-study.pptx")

# --- DECK 3: DEEP-DIVE ---
def build_deck_3():
    prs = create_deck("evolith-sdlc-technical-deepdive.pptx", "Evolith SDLC: Deep-Dive Técnico", "Diagramas de Arquitectura y Pipeline")
    
    # Slide: CI/CD Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_theme(slide)
    add_title(slide, "Línea de Ensamblaje Autónoma (F4 y F5)")
    
    add_conceptual_box(slide, 0.5, 3.0, 1.5, 1.0, "Git Push", WHITE_TEXT)
    add_arrow(slide, 2.1, 3.3, 0.4, 0.4)
    add_conceptual_box(slide, 2.6, 3.0, 1.5, 1.0, "Build", WHITE_TEXT)
    add_arrow(slide, 4.2, 3.3, 0.4, 0.4)
    add_conceptual_box(slide, 4.7, 2.5, 1.5, 2.0, "SAST Scan\nQuality Gate", DARK_BG, GOLD_ACCENT)
    add_arrow(slide, 6.3, 3.3, 0.4, 0.4)
    add_conceptual_box(slide, 6.8, 3.0, 1.5, 1.0, "Deploy Staging", WHITE_TEXT)
    
    # Slide: Blue/Green
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_theme(slide)
    add_title(slide, "Despliegue Sin Interrupciones (F6)")
    
    add_conceptual_box(slide, 1.0, 3.5, 2.0, 1.0, "Load Balancer", GOLD_ACCENT)
    add_conceptual_box(slide, 5.0, 2.0, 3.0, 1.5, "Entorno BLUE (V1)\n(Drenando)", GRAY_TEXT)
    add_conceptual_box(slide, 5.0, 4.5, 3.0, 1.5, "Entorno GREEN (V2)\n(Tráfico Activo)", DARK_BG, GOLD_ACCENT)
    
    save_deck(prs, "evolith-sdlc-technical-deepdive.pptx")

# --- EXCEL HELPERS ---
def format_sheet(sheet, headers, data):
    sheet.append(headers)
    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="121212", end_color="121212", fill_type="solid") # Onyx/Black
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = Border(left=Side(style='thin', color="CAC5B2"), right=Side(style='thin', color="CAC5B2"), top=Side(style='thin', color="CAC5B2"), bottom=Side(style='thin', color="CAC5B2"))
    for idx, r in enumerate(data):
        sheet.append(r)
        fill_color = "2A2A2A" if idx % 2 == 0 else "1A1A1A" # Dark theme rows
        for cell in sheet[sheet.max_row]:
            cell.font = Font(color="FFFFFF")
            cell.fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
            cell.alignment = Alignment(wrapText=True, vertical="center")
            cell.border = Border(left=Side(style='thin', color="CAC5B2"), right=Side(style='thin', color="CAC5B2"), top=Side(style='thin', color="CAC5B2"), bottom=Side(style='thin', color="CAC5B2"))
    for col in sheet.columns:
        max_length = 0
        col_letter = col[0].column_letter
        for cell in col:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        sheet.column_dimensions[col_letter].width = min(max_length + 5, 45)

def add_portada_excel(wb, title):
    ws = wb.active
    ws.title = "00_Portada"
    ws.sheet_properties.tabColor = "CAC5B2"
    ws.append(["EVOLITH SDLC - " + title])
    ws.append(["Versión 1.0"])
    ws.append(["Nivel de Confidencialidad: Uso Interno"])
    ws.append(["Instrucciones: Llenar cada pestaña conforme se avanza en la fase correspondiente."])
    
    # Dark mode portada
    for row in ws.iter_rows(min_row=1, max_row=10, min_col=1, max_col=5):
        for cell in row:
            cell.fill = PatternFill(start_color="121212", end_color="121212", fill_type="solid")
            cell.font = Font(color="FFFFFF")
            
    ws['A1'].font = Font(bold=True, size=16, color="CAC5B2") # Gold title

def build_xlsxs():
    # 1. Master Workbook
    wb1 = Workbook()
    add_portada_excel(wb1, "Workbook Maestro Integrador")
    ws = wb1.create_sheet("Dashboard_Proyecto")
    format_sheet(ws, ["ID_Indicador", "Nombre Métrica", "Valor Actual", "Objetivo", "Estado", "Tendencia"], [
        ["IND-01", "Cumplimiento Cronograma", "95%", "> 90%", "Verde", "Estable"],
        ["IND-02", "Cobertura de Pruebas", "82%", "> 80%", "Verde", "Al alza"],
        ["IND-03", "Defectos Abiertos Severidad 1", "2", "0", "Rojo", "Estable"]
    ])
    wb1.save(os.path.join(out_dir, "Evolith_Master_Workbook.xlsx"))

    # 2. Ideacion
    wb2 = Workbook()
    add_portada_excel(wb2, "Fase 1: Ideación")
    ws = wb2.create_sheet("Business_Case")
    format_sheet(ws, ["Campo", "Valor Propuesto"], [
        ["Nombre Iniciativa", "Portal de Autoservicio Industrial"],
        ["Problema a Resolver", "Alta carga operativa telefónica."],
        ["Beneficio Cuantificable", "Reducción 40% llamadas L1."],
        ["Costo Estimado (CapEx)", "$150,000 USD"],
        ["Decisión Comité", "Aprobado"]
    ])
    wb2.save(os.path.join(out_dir, "Evolith_Workbook_F1_Ideacion.xlsx"))

    # 3. Analisis
    wb3 = Workbook()
    add_portada_excel(wb3, "Fase 2: Análisis")
    ws = wb3.create_sheet("Especificacion_Requerimientos")
    format_sheet(ws, ["Módulo", "ID_Req", "Tipo", "Descripción Funcional", "Prioridad", "Criterios Aceptación", "Estado"], [
        ["Autenticación", "REQ-AUTH-01", "Funcional", "Login industrial con SSO.", "Must Have", "Redirección exitosa", "Aprobado"]
    ])
    wb3.save(os.path.join(out_dir, "Evolith_Workbook_F2_Analisis.xlsx"))

    # 4. Diseño
    wb4 = Workbook()
    add_portada_excel(wb4, "Fase 3: Diseño")
    ws = wb4.create_sheet("Blueprint")
    format_sheet(ws, ["Componente", "Descripción", "Interfaz / API", "Atributo Calidad"], [
        ["BFF", "Gateway Web", "REST / GraphQL", "Alta disponibilidad"]
    ])
    wb4.save(os.path.join(out_dir, "Evolith_Workbook_F3_Diseno.xlsx"))

    # 5. Construccion
    wb5 = Workbook()
    add_portada_excel(wb5, "Fase 4: Construcción")
    ws = wb5.create_sheet("Pipeline_CI")
    format_sheet(ws, ["Etapa", "Herramienta", "Checkpoint Calidad", "Estado"], [
        ["Static Analysis", "SonarQube", "Cobertura > 80%", "Bloqueante"]
    ])
    wb5.save(os.path.join(out_dir, "Evolith_Workbook_F4_Construccion.xlsx"))

    # 6. Pruebas
    wb6 = Workbook()
    add_portada_excel(wb6, "Fase 5: Pruebas")
    ws = wb6.create_sheet("Casos_Prueba")
    format_sheet(ws, ["ID_Caso", "Requisito", "Pasos", "Resultado Esperado", "Estado"], [
        ["CP-001", "REQ-AUTH-01", "Entrar web, Poner creds", "Redirección a Home", "Passed"]
    ])
    wb6.save(os.path.join(out_dir, "Evolith_Workbook_F5_Pruebas.xlsx"))

    # 7. Despliegue
    wb7 = Workbook()
    add_portada_excel(wb7, "Fase 6: Despliegue")
    ws = wb7.create_sheet("Runbook")
    format_sheet(ws, ["Paso", "Tipo", "Responsable", "Rollback"], [
        ["1. Backup DB", "Automático", "DBA", "Restaurar Snapshot"]
    ])
    wb7.save(os.path.join(out_dir, "Evolith_Workbook_F6_Despliegue.xlsx"))

    # 8. Operacion
    wb8 = Workbook()
    add_portada_excel(wb8, "Fase 7: Operación")
    ws = wb8.create_sheet("Dashboard")
    format_sheet(ws, ["Métrica", "Valor Actual", "Umbral SLA", "Estado"], [
        ["Uptime API", "99.95%", "> 99.9%", "Saludable"]
    ])
    wb8.save(os.path.join(out_dir, "Evolith_Workbook_F7_Operacion.xlsx"))

    # 9. Retiro
    wb9 = Workbook()
    add_portada_excel(wb9, "Fase 8: Retiro")
    ws = wb9.create_sheet("Plan_Retiro")
    format_sheet(ws, ["Sistema Legacy", "Fecha Sunset", "Estrategia Migración"], [
        ["Portal Legacy v2", "2026-12-31", "Exportación CSV a Postgres."]
    ])
    wb9.save(os.path.join(out_dir, "Evolith_Workbook_F8_Retiro.xlsx"))

if __name__ == "__main__":
    print("Generando Kit Completo de 12 Archivos (IDENTIDAD VISUAL BEYONDNET - OSCURO/DORADO)...")
    build_deck_1()
    build_deck_2()
    build_deck_3()
    build_xlsxs()
    print("Todo generado exitosamente en assets/")
